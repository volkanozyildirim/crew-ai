"""Sprint report generator — builds a FLO-branded .pptx for a sprint.

Groups a sprint's work items by their parent (Epic/Feature = "project"), rolls up
a status per group, summarizes the work done (reusing already-stored data:
MySQL pipeline step outputs -> vector-store /repo-decisions -> Azure WI text), and
injects it into the "Sprint Özeti — Neler Yaptık" slides of a cloned template deck.

Everything user-facing (slide titles, status suffixes, AI bullets) is Turkish;
code + LLM instructions are English (repo convention).

Feature is env-gated by CREW_SPRINT_REPORT. Template + output live outside git.
"""
from __future__ import annotations

import copy
import json
import logging
import os
import re
from dataclasses import dataclass, field
from datetime import datetime

log = logging.getLogger("agile_sdlc_crew.sprint_report")

# ------------------------------------------------------------------ config

DONE_STATES = {"done", "closed", "resolved", "completed", "ready for production"}


def _env(name: str, default: str = "") -> str:
    return os.environ.get(name, default)


def template_path() -> str:
    p = _env("CREW_SPRINT_REPORT_TEMPLATE") or os.path.expanduser(
        "~/.crew_repos/sprint_report_template.pptx"
    )
    return os.path.expanduser(p)


def output_dir() -> str:
    p = _env("CREW_SPRINT_REPORT_OUT_DIR") or os.path.expanduser(
        "~/.crew_repos/sprint_reports"
    )
    p = os.path.expanduser(p)
    os.makedirs(p, exist_ok=True)
    return p


def _groups_per_slide() -> int:
    # görseller açıkken her parent kendi slaytında (1); aksi halde env/3
    if _visuals_enabled():
        default = "1"
    else:
        default = "3"
    try:
        return max(1, int(_env("CREW_SPRINT_REPORT_GROUPS_PER_SLIDE", default)))
    except ValueError:
        return 1 if _visuals_enabled() else 3


def _ai_enabled() -> bool:
    return _env("CREW_SPRINT_REPORT_AI_SUMMARY", "1") not in ("0", "false", "False", "")


def _visuals_enabled() -> bool:
    return _env("CREW_SPRINT_REPORT_VISUALS", "1") not in ("0", "false", "False", "")


def _azure_charts_enabled() -> bool:
    """Burndown + velocity grafiklerini Azure Analytics'ten otomatik çiz."""
    return _env("CREW_SPRINT_REPORT_AZURE_CHARTS", "1") not in ("0", "false", "False", "")


def _velocity_sprints() -> int:
    try:
        return max(1, int(_env("CREW_SPRINT_REPORT_VELOCITY_SPRINTS", "6")))
    except ValueError:
        return 6


def _min_pct() -> float:
    """Ayrı slayt+rozet için eşik: sprint toplam SP'sinin bu %'sinden FAZLA efor
    alan parent'lar ayrı slayta; altındakiler 'Diğer Çalışmalar' özetine gider."""
    try:
        return float(_env("CREW_SPRINT_REPORT_MIN_PCT", "5"))
    except ValueError:
        return 5.0


def _digest_per_slide() -> int:
    """'Diğer Çalışmalar' slaytında satır (proje) sayısı."""
    try:
        return max(1, int(_env("CREW_SPRINT_REPORT_DIGEST_PER_SLIDE", "6")))
    except ValueError:
        return 6


def icon_dir() -> str:
    p = _env("CREW_SPRINT_REPORT_ICON_DIR") or os.path.expanduser(
        "~/.crew_repos/sprint_report_icons"
    )
    p = os.path.expanduser(p)
    os.makedirs(p, exist_ok=True)
    return p


def team_image_path() -> str:
    """Ekip slaytında kullanılacak takım görseli (varsa eski görselin yerine geçer)."""
    return os.path.expanduser(
        _env("CREW_SPRINT_REPORT_TEAM_IMAGE")
        or os.path.expanduser("~/.crew_repos/sprint_report_team.png")
    )


# Domain sınıfları: anahtar -> (Türkçe etiket, emoji-fallback, ikon konsepti, anahtar kelimeler)
FLO_ORANGE = "#E97132"
NAVY = "#0E2841"

DOMAIN_META = {
    "kargo": ("KARGO & TESLİMAT", "🚚", "a delivery truck",
              ["kargo", "teslimat", "dhl", "cargox", "kargoist", "5post", "fivepost",
               "dalli", "dağıtım", "dagitim", "droppoint", "shipment", "gönderi",
               "sevkiyat", "randevulu", "flokargom", "flo kargom"]),
    "stok": ("STOK & DEPO", "📦", "a warehouse box / inventory stack",
             ["stok", "stock", "warehouse", "depo", "alokasyon", "allocation",
              "envanter", "transit", "sku"]),
    "entegrasyon": ("ENTEGRASYON", "🔌", "two connected plugs / API integration",
                    ["entegrasyon", "integration", "servis", "api", "atol", "ozon",
                     "connector", "couchbase", "kafka", "orkestra"]),
    "iade": ("İADE", "↩️", "a return arrow with a package",
             ["iade", "return", "returnoffices"]),
    "odeme": ("ÖDEME & FİYAT", "💳", "a credit card / price tag",
              ["ödeme", "odeme", "payment", "fiyat", "ücret", "ucret", "limit",
               "fatura", "shipping_amount", "amount"]),
    "kanal": ("KANAL & PANEL", "🖥️", "a dashboard panel screen",
              ["kanal", "channel", "panel", "yayın", "yayin", "marketplace", "butterfly"]),
    "siparis": ("SİPARİŞ", "🧾", "an order receipt",
                ["sipariş", "siparis", "order", "crm", "setshipping"]),
    "rapor": ("RAPOR & İZLEME", "📊", "a bar chart / monitoring dashboard",
              ["rapor", "izleme", "filtre", "ekran", "dashboard", "log", "raporlan"]),
    "diger": ("DİĞER", "🗂️", "a folder of documents", []),
}


def _classify_domain(text: str) -> str:
    """Parent+iş metninde en çok anahtar kelime geçen domain'i seç (baskın alan).
    Hiç eşleşme yoksa 'diger'. Eşitlikte priority sırası kazanır."""
    t = (text or "").lower()
    # eşitlik bozucu öncelik (lojistik ekibi için çekirdek alanlar önce)
    priority = ["kargo", "stok", "entegrasyon", "kanal", "siparis", "iade", "odeme", "rapor"]
    scores = {}
    for key in priority:
        _, _, _, kws = DOMAIN_META[key]
        scores[key] = sum(t.count(kw) for kw in kws)
    best = max(priority, key=lambda k: (scores[k], -priority.index(k)))
    return best if scores[best] > 0 else "diger"


# ------------------------------------------------------------------ data model


@dataclass
class SprintItem:
    id: int
    title: str
    type: str = ""
    state: str = ""
    assignee: str = ""
    story_points: float = 0.0
    parent_id: int | None = None
    parent_title: str = ""
    parent_type: str = ""


@dataclass
class ParentGroup:
    parent_id: int | None
    parent_title: str
    parent_type: str
    items: list[SprintItem] = field(default_factory=list)

    @property
    def status(self) -> str:
        """Türkçe rolled-up durum: tüm çocuklar bittiyse Tamamlandı, değilse Devam Ediyor."""
        if self.items and all(i.state.strip().lower() in DONE_STATES for i in self.items):
            return "Tamamlandı"
        return "Devam Ediyor"

    @property
    def sp_total(self) -> float:
        """Gruptaki (parent altındaki) toplam tüketilen story point."""
        return sum(i.story_points or 0.0 for i in self.items)


@dataclass
class GroupSummary:
    title: str
    status: str
    bullets: list[str]
    sp: float = 0.0          # bu parent altında tüketilen toplam SP
    sp_pct: float = 0.0      # sprint toplam SP'si içindeki oran (%)
    domain: str = "diger"    # görsel domain anahtarı (DOMAIN_META)


# ------------------------------------------------------------------ data collection


def collect_sprint_items(client, iteration_path: str) -> list[SprintItem]:
    """Sprint'teki work item'ları parent'larıyla birlikte topla."""
    raw = client.get_iteration_work_items(iteration_path)
    items = [
        SprintItem(
            id=int(r["id"]),
            title=r.get("title", "") or "",
            type=r.get("type", "") or "",
            state=r.get("state", "") or "",
            assignee=r.get("assignedTo", "") or "",
            story_points=float(r.get("storyPoints") or 0.0),
        )
        for r in raw
        if r.get("id") is not None
    ]
    if not items:
        return items
    try:
        parents = client.get_work_item_parents([i.id for i in items])
    except Exception as e:  # pragma: no cover - network
        log.warning(f"  Parent çözümlemesi başarısız, gruplama parent'sız: {e}")
        parents = {}
    for it in items:
        p = parents.get(it.id)
        if p:
            it.parent_id = p.get("parent_id")
            it.parent_title = p.get("parent_title", "") or ""
            it.parent_type = p.get("parent_type", "") or ""
    return items


def group_by_parent(items: list[SprintItem]) -> list[ParentGroup]:
    """Parent'a göre grupla; parent'ı olmayanlar 'Diğer / Bağımsız İşler' altında."""
    groups: dict[object, ParentGroup] = {}
    for it in items:
        key = it.parent_id if it.parent_id else "__none__"
        if key not in groups:
            if key == "__none__":
                groups[key] = ParentGroup(None, "Diğer / Bağımsız İşler", "")
            else:
                groups[key] = ParentGroup(
                    it.parent_id, it.parent_title or f"#{it.parent_id}", it.parent_type
                )
        groups[key].items.append(it)
    # tüketilen story point'e göre AZALAN; parent'sız 'Diğer' grubu en sonda
    ordered = sorted(
        groups.values(),
        key=lambda g: (g.parent_id is None, -g.sp_total, g.parent_title.lower()),
    )
    return ordered


def _stored_context_for_item(item: SprintItem, db, repo_decisions: dict, client) -> str:
    """'Ne yapıldı' bağlamını ZATEN DEPOLANMIŞ kaynaklardan çek (yeniden keşif yok):
    1) MySQL pipeline çıktısı  2) vector-store /repo-decisions  3) Azure WI metni (fallback).
    """
    wi = str(item.id)
    # 1) MySQL pipeline outputs (en zengin)
    if db is not None:
        for step in ("technical_design_task", "requirements_analysis_task"):
            try:
                out = db.get_cached_step_output(step, wi)
            except Exception:
                out = None
            if out and out.strip():
                return out.strip()[:1500]
    # 2) vector-store /repo-decisions (backfill'lenmiş done WI'ler)
    rec = repo_decisions.get(wi)
    if rec and rec.get("content"):
        return str(rec["content"])[:1500]
    # 3) Azure WI açıklaması / AC (fallback)
    if client is not None:
        try:
            full = client.get_work_item(item.id)
            f = full.get("fields", {}) if isinstance(full, dict) else {}
            parts = [
                _strip_html(f.get("System.Description", "")),
                _strip_html(f.get("Microsoft.VSTS.Common.AcceptanceCriteria", "")),
            ]
            txt = "\n".join(p for p in parts if p).strip()
            if txt:
                return txt[:1500]
        except Exception:
            pass
    return ""


def _strip_html(s: str) -> str:
    if not s:
        return ""
    s = re.sub(r"<br\s*/?>", "\n", s, flags=re.IGNORECASE)
    s = re.sub(r"</(p|div|li)>", "\n", s, flags=re.IGNORECASE)
    s = re.sub(r"<[^>]+>", "", s)
    s = s.replace("&nbsp;", " ").replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    return re.sub(r"[ \t]+", " ", s).strip()


# ------------------------------------------------------------------ summarization


def summarize_group(group: ParentGroup, contexts: dict[int, str]) -> GroupSummary:
    """Grup için başlık + durum + Türkçe maddeler üret (AI açıksa AI, değilse WI başlıkları)."""
    title = group.parent_title
    status = group.status
    sp = group.sp_total
    # domain: parent başlığı + iş başlıkları üzerinden sınıfla
    clf_text = title + " " + " ".join(it.title for it in group.items)
    domain = _classify_domain(clf_text)
    if _ai_enabled():
        bullets = _ai_bullets(group, contexts)
        if bullets:
            return GroupSummary(title, status, bullets, sp=sp, domain=domain)
    # fallback: her WI başlığı bir madde
    bullets = [f"{it.title}" for it in group.items][:8]
    return GroupSummary(title, status, bullets, sp=sp, domain=domain)


def _ai_bullets(group: ParentGroup, contexts: dict[int, str]) -> list[str]:
    """claude_cli ile SADECE toplanmış bağlamdan 2-4 Türkçe madde. Hata olursa [] döner."""
    try:
        from agile_sdlc_crew.llm import build_for_agent
    except Exception as e:
        log.warning(f"  AI özet atlandı (LLM import): {e}")
        return []
    lines = []
    for it in group.items:
        ctx = contexts.get(it.id, "")
        lines.append(f"- WI #{it.id} [{it.state}] {it.title}\n  {ctx[:800]}")
    context_block = "\n".join(lines)[:6000]
    prompt = (
        "You are writing a sprint-review slide for FLO GROUP. Below is a group of "
        "work items that belong to the same parent project, with already-collected "
        "context about what was implemented.\n\n"
        f"Parent project: {group.parent_title}\n"
        f"Work items:\n{context_block}\n\n"
        "Write a concise summary of what was accomplished for this project, as 2-4 "
        "short bullet points. Rules:\n"
        "- Write the bullets in TURKISH.\n"
        "- Each bullet: one clear outcome, past tense, no work-item numbers.\n"
        "- Do NOT invent anything not supported by the context; if context is thin, "
        "summarize from the work item titles.\n"
        'Return ONLY strict JSON: {"bullets": ["...", "..."]}'
    )
    try:
        llm = build_for_agent("business_analyst")
        resp = llm.call(prompt)
    except Exception as e:
        log.warning(f"  AI özet çağrısı başarısız (grup '{group.parent_title}'): {e}")
        return []
    return _parse_bullets(resp)


def _parse_bullets(resp: str) -> list[str]:
    if not resp:
        return []
    m = re.search(r"\{.*\}", resp, re.DOTALL)
    if m:
        try:
            data = json.loads(m.group(0))
            bl = data.get("bullets")
            if isinstance(bl, list):
                return [str(b).strip() for b in bl if str(b).strip()][:4]
        except Exception:
            pass
    # fallback: satır bazlı
    out = []
    for ln in resp.splitlines():
        ln = ln.strip().lstrip("-•*").strip()
        if ln and not ln.startswith("{") and not ln.startswith("}"):
            out.append(ln)
    return out[:4]


# ------------------------------------------------------------------ pptx assembly

_NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _clone_slide(prs, src):
    """src slaydının birebir kopyasını (görsel ilişkileriyle) deste sonuna ekler."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    layout = src.slide_layout
    new = prs.slides.add_slide(layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    remap = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        if rel.is_external:
            new_rId = new.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
        else:
            new_rId = new.part.relate_to(rel.target_part, rel.reltype)
        remap[rId] = new_rId
    for sh in src.shapes:
        el = copy.deepcopy(sh._element)
        for node in el.iter():
            for attr in (_NS_R + "embed", _NS_R + "link"):
                if attr in node.attrib and node.attrib[attr] in remap:
                    node.attrib[attr] = remap[node.attrib[attr]]
        new.shapes._spTree.append(el)
    return new


def _slide_text(slide) -> str:
    parts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            parts.append(sh.text_frame.text)
    return " ".join(parts)


def _find_shape(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def _content_textbox(slide):
    """İçerik slaytındaki gövde metin kutusu ('TextBox 1' ya da en geniş metin kutusu)."""
    box = _find_shape(slide, "TextBox 1")
    if box is not None:
        return box
    # fallback: başlık/lorem dışındaki en büyük text box
    best = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.name.startswith("TextBox"):
            if best is None or (sh.width or 0) * (sh.height or 0) > (best.width or 0) * (best.height or 0):
                best = sh
    return best


def _ai_svg_icon(concept: str) -> str:
    """claude_cli ile kurumsal, tek renk flat-line SVG ikon markup'ı üret."""
    try:
        from agile_sdlc_crew.llm import build_for_agent

        llm = build_for_agent("business_analyst")
    except Exception as e:
        log.warning(f"  SVG ikon LLM alınamadı: {e}")
        return ""
    prompt = (
        f"Generate a clean, professional, CORPORATE flat line icon as SVG for: '{concept}'.\n"
        "Strict requirements:\n"
        "- Output ONLY the <svg>...</svg> markup, no explanation, no code fences.\n"
        "- Attributes: xmlns='http://www.w3.org/2000/svg', viewBox='0 0 24 24', width='24', height='24'.\n"
        f"- Outline style only: stroke='{FLO_ORANGE}', stroke-width='1.6', fill='none', "
        "stroke-linecap='round', stroke-linejoin='round'.\n"
        "- No text, no background rectangle, no gradients, no drop shadows.\n"
        "- Minimal, recognizable, balanced, ~2 units padding inside the viewBox."
    )
    try:
        resp = llm.call(prompt)
    except Exception as e:
        log.warning(f"  SVG ikon çağrısı başarısız ({concept}): {e}")
        return ""
    m = re.search(r"<svg\b.*?</svg>", resp or "", re.DOTALL | re.IGNORECASE)
    return m.group(0) if m else ""


def _domain_icon_png(domain: str) -> str | None:
    """Domain PNG ikon yolu; yoksa claude SVG + cairosvg ile üretip cache'le (bir kez)."""
    meta = DOMAIN_META.get(domain, DOMAIN_META["diger"])
    png = os.path.join(icon_dir(), f"{domain}.png")
    if os.path.exists(png):
        return png
    svg = _ai_svg_icon(meta[2])
    if not svg:
        return None
    # modelin döndürdüğü rengi FLO turuncuya sabitle
    svg = re.sub(r'stroke="#[0-9A-Fa-f]{3,6}"', f'stroke="{FLO_ORANGE}"', svg)
    try:
        import cairosvg

        cairosvg.svg2png(
            bytestring=svg.encode("utf-8"), write_to=png,
            output_width=320, output_height=320,
        )
    except Exception as e:
        log.warning(f"  İkon rasterize edilemedi ({domain}): {e}")
        return None
    return png


def _add_domain_badge(slide, summary: GroupSummary):
    """Üst-sağ köşeye domain rozeti: beyaz kart + FLO turuncu ikon + Türkçe etiket."""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    meta = DOMAIN_META.get(summary.domain, DOMAIN_META["diger"])
    label = meta[0]
    left, top, w, h = Inches(10.45), Inches(1.05), Inches(2.55), Inches(1.95)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xE9, 0x71, 0x32)
    card.line.width = Pt(1.25)
    try:
        card.shadow.inherit = False
    except Exception:
        pass

    icon = _domain_icon_png(summary.domain)
    if icon:
        iw = Inches(1.05)
        slide.shapes.add_picture(icon, left + (w - iw) // 2, top + Inches(0.2), height=iw)
    else:
        eb = slide.shapes.add_textbox(left, top + Inches(0.12), w, Inches(1.05))
        ep = eb.text_frame.paragraphs[0]
        ep.alignment = PP_ALIGN.CENTER
        er = ep.add_run()
        er.text = meta[1]
        er.font.size = Pt(44)

    lb = slide.shapes.add_textbox(left, top + Inches(1.34), w, Inches(0.5))
    lb.text_frame.word_wrap = True
    lp = lb.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = label
    lr.font.size = Pt(12)
    lr.font.bold = True
    lr.font.color.rgb = RGBColor(0x0E, 0x28, 0x41)


def _fill_content_slide(slide, summaries: list[GroupSummary], start_no: int):
    """İçerik slaytını doldur: 'TextBox 1' içine numaralı proje başlıkları + maddeler."""
    from pptx.util import Inches, Pt

    # ortadaki ayraç çizgilerini kaldır (sabit konumlu, değişken içerikle hizasız)
    for sh in list(slide.shapes):
        try:
            from pptx.util import Emu

            if sh.shape_type == 9 and Emu(sh.top).inches > 1.3:  # LINE, gövdedeki
                sh._element.getparent().remove(sh._element)
        except Exception:
            pass

    box = _content_textbox(slide)
    if box is None:
        log.warning("  İçerik metin kutusu bulunamadı; slayt atlandı")
        return
    # görseller açık ve slaytta tek parent -> sağda rozete yer bırak (kutuyu daralt)
    do_badge = _visuals_enabled() and len(summaries) == 1
    # tüm içerik slaytlarında AYNI hizalama (şablonun kutu top'ları farklı olabiliyor)
    box.left = Inches(1.5)
    box.top = Inches(1.3)
    box.width = Inches(8.6) if do_badge else Inches(11.6)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    no = start_no
    for s in summaries:
        # başlık paragrafı
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            # gruplar arası boş satır
            sp = tf.add_paragraph()
            sp.level = 0
            p = tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = f"{no}. {s.title} – {s.status}"
        r.font.size = Pt(16)
        r.font.bold = True
        # SP + sprint toplamı içindeki oran (bold değil, biraz daha küçük)
        if s.sp:
            r2 = p.add_run()
            r2.text = f"   ({s.sp:g} SP · %{int(round(s.sp_pct))})"
            r2.font.size = Pt(13)
            r2.font.bold = False
        # maddeler
        for b in s.bullets:
            bp = tf.add_paragraph()
            bp.level = 1
            bp.space_before = Pt(4)
            bp.space_after = Pt(4)
            bp.line_spacing = 1.1
            br = bp.add_run()
            br.text = f"•  {b}"  # görünür madde imi
            br.font.size = Pt(14)
            br.font.bold = False
        no += 1

    if do_badge:
        _add_domain_badge(slide, summaries[0])


def _set_slide_title(slide, text: str):
    """İçerik slaytının başlık kutusunu (Sprint Özeti - Neler Yaptık) değiştir."""
    for sh in slide.shapes:
        if sh.has_text_frame and "Neler Yaptık" in sh.text_frame.text:
            p = sh.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = text
            return


def _fill_digest_slide(slide, summaries: list[GroupSummary]):
    """'Diğer Çalışmalar' slaytı: her proje = başlık satırı (SP·%) + 1 cümle özet."""
    from pptx.util import Emu, Inches, Pt

    # ortadaki ayraç çizgilerini kaldır
    for sh in list(slide.shapes):
        try:
            if sh.shape_type == 9 and Emu(sh.top).inches > 1.3:
                sh._element.getparent().remove(sh._element)
        except Exception:
            pass

    box = _content_textbox(slide)
    if box is None:
        return
    box.left = Inches(1.5)  # içerik slaytlarıyla aynı hiza
    box.top = Inches(1.3)
    box.width = Inches(11.6)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for s in summaries:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            sp = tf.add_paragraph()  # gruplar arası küçük boşluk
            sp.level = 0
            p = tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = f"{s.title} – {s.status}"
        r.font.size = Pt(14)
        r.font.bold = True
        if s.sp:
            r2 = p.add_run()
            r2.text = f"   ({s.sp:g} SP · %{int(round(s.sp_pct))})"
            r2.font.size = Pt(11)
            r2.font.bold = False
        # tek cümlelik özet (ilk madde)
        line = s.bullets[0] if s.bullets else ""
        if line:
            bp = tf.add_paragraph()
            bp.level = 1
            br = bp.add_run()
            br.text = line
            br.font.size = Pt(12)
            br.font.bold = False


def _add_burndown_chart(slide, data: dict):
    """Burndown: Kalan = dolu mavi alan; Toplam Kapsam = turuncu çizgi; İdeal = gri çizgi.
    (Tek AREA chart; çizgi serilerinin dolgusu kapatılıp yalnız üst kenar çizgisi bırakılır.)"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    cd = CategoryChartData()
    cd.categories = data["labels"]
    cd.add_series("Kalan (SP)", data["remaining"])
    cd.add_series("Toplam Kapsam", data["total"])
    cd.add_series("İdeal", data["ideal"])
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.AREA, Inches(1.4), Inches(1.2), Inches(10.5), Inches(4.95), cd
    )
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False

    s0, s1, s2 = ch.series[0], ch.series[1], ch.series[2]
    # Kalan: dolu mavi alan
    s0.format.fill.solid()
    s0.format.fill.fore_color.rgb = RGBColor(0x41, 0x7A, 0xD6)
    s0.format.line.color.rgb = RGBColor(0x41, 0x7A, 0xD6)
    # Toplam Kapsam: dolgu yok -> yalnız turuncu üst kenar çizgisi
    s1.format.fill.background()
    s1.format.line.color.rgb = RGBColor(0xE9, 0x71, 0x32)
    s1.format.line.width = Pt(2.25)
    # İdeal: dolgu yok -> gri çizgi
    s2.format.fill.background()
    s2.format.line.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)
    s2.format.line.width = Pt(1.5)
    return ch


def _add_velocity_chart(slide, data: dict):
    """Değerlendirme slaytına son-N sprint clustered column chart (Planlanan/Toplam/Tamamlanan)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    cd = CategoryChartData()
    cd.categories = data["sprints"]
    cd.add_series("Planlanan Efor", data["planned"])
    cd.add_series("Toplam Efor", data["total"])
    cd.add_series("Tamamlanan Efor", data["done"])
    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.2), Inches(1.2), Inches(10.9), Inches(4.95), cd
    )
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    for i, rgb in enumerate([(0x4C, 0x96, 0xF0), (0xA0, 0x50, 0xDC), (0x5A, 0xC8, 0x78)]):
        ser = ch.plots[0].series[i]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = RGBColor(*rgb)
    # x ekseni (sprint adları) etiketlerini küçült
    try:
        ch.category_axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass
    return ch


def _add_azure_charts(prs, burndown: dict | None, velocity: dict | None):
    """Metrik slaytlarını Azure verisiyle çiz (burndown + velocity)."""
    for slide in prs.slides:
        t = _slide_text(slide)
        if burndown and "Burndown" in t:
            try:
                _add_burndown_chart(slide, burndown)
            except Exception as e:
                log.warning(f"  Burndown grafiği eklenemedi: {e}")
        elif velocity and ("değerlendirme" in t or "Son 3 ayl" in t):
            try:
                _add_velocity_chart(slide, velocity)
            except Exception as e:
                log.warning(f"  Velocity grafiği eklenemedi: {e}")


def _set_team_image(prs, image_path: str):
    """Ekip slaytındaki büyük takım görselini verilen görselle değiştir (aynı konum/boyut)."""
    from pptx.util import Inches

    if not image_path or not os.path.exists(image_path):
        return
    for slide in prs.slides:
        t = _slide_text(slide)
        if "Ekip" not in t or "Neler Yaptık" in t or "Diğer Çalışmalar" in t or "Teşekkür" in t:
            continue
        box = None
        for sh in list(slide.shapes):
            if sh.shape_type == 13 and (sh.width or 0) >= Inches(3):  # büyük ekip görseli
                box = (sh.left, sh.top, sh.width, sh.height)
                sh._element.getparent().remove(sh._element)
        if box is None:
            box = (Inches(2.83), Inches(1.06), Inches(7.68), Inches(5.86))
        slide.shapes.add_picture(image_path, box[0], box[1], width=box[2], height=box[3])
        return


def _set_title_slide(prs, title_line: str, sprint_name: str, date_range: str):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    for slide in prs.slides:
        box = _find_shape(slide, "Metin kutusu 2")
        if box is None:
            continue
        tf = box.text_frame
        tf.clear()
        specs = [(title_line, 32, True), (sprint_name, 32, False), (date_range, 20, False)]
        specs = [s for s in specs if s[0]]  # boş satırları (ör. tarihsiz sprint) atla
        for i, (txt, sz, bold) in enumerate(specs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(sz)
            r.font.bold = bold
        return  # yalnızca ilk (başlık) slayt


def _strip_metric_slides(prs):
    """Metrik slaytlarındaki eski ekran görüntülerini/tabloları kaldır (başlıklı boş bırak)."""
    from pptx.util import Inches

    for slide in prs.slides:
        txt = _slide_text(slide)
        if "Burndown" not in txt and "değerlendirme" not in txt and "Son 3 ayl" not in txt:
            continue
        for sh in list(slide.shapes):
            try:
                is_pic = sh.shape_type == 13
                is_tbl = sh.shape_type == 19
                wide = (sh.width or 0) >= Inches(3)
                if (is_pic and wide) or is_tbl:
                    sh._element.getparent().remove(sh._element)
            except Exception:
                pass


def _renumber_pages(prs):
    from pptx.util import Pt

    for idx, slide in enumerate(prs.slides):
        box = _find_shape(slide, "9")
        if box is not None and box.has_text_frame:
            try:
                p = box.text_frame.paragraphs[0]
                p.text = str(idx + 1)  # sayfa no
                for r in p.runs:
                    r.font.size = Pt(12)
            except Exception:
                pass


def _reorder_slides(prs):
    """Slaytları kanonik sıraya diz: başlık, burndown, içerik*, 3-aylık, ekip, teşekkürler."""

    def rank(slide):
        t = _slide_text(slide)
        if _find_shape(slide, "Metin kutusu 2") is not None and "Teşekkürler" not in t:
            return 0
        if "Burndown" in t:
            return 1
        if "Neler Yaptık" in t or "Diğer Çalışmalar" in t:
            return 2
        if "değerlendirme" in t or "Son 3 ayl" in t:
            return 3
        if "Teşekkürler" in t:
            return 5
        if "Ekip" in t:
            return 4
        return 6

    sldIdLst = prs.slides._sldIdLst
    id_els = list(sldIdLst)
    slides = list(prs.slides)
    order = sorted(range(len(slides)), key=lambda i: rank(slides[i]))  # kararlı sıralama
    for el in id_els:
        sldIdLst.remove(el)
    for i in order:
        sldIdLst.append(id_els[i])


def build_pptx(
    title_line: str,
    sprint_name: str,
    date_range: str,
    summaries: list[GroupSummary],
    tmpl_path: str,
    out_path: str,
    burndown: dict | None = None,
    velocity: dict | None = None,
) -> str:
    from pptx import Presentation

    prs = Presentation(tmpl_path)

    _set_title_slide(prs, title_line, sprint_name, date_range)
    _strip_metric_slides(prs)
    _add_azure_charts(prs, burndown, velocity)

    # şablondaki mevcut içerik slaytları (slides ile sldIdLst aynı sırada -> paralel indeks)
    sldIdLst = prs.slides._sldIdLst
    all_slides = list(prs.slides)
    all_els = list(sldIdLst)
    el_of = {i: all_els[i] for i in range(len(all_slides))}
    content_idx = [i for i, s in enumerate(all_slides) if "Neler Yaptık" in _slide_text(s)]
    if not content_idx:
        raise RuntimeError("Şablonda 'Neler Yaptık' içerik slaytı bulunamadı")
    content_slides = [all_slides[i] for i in content_idx]
    src = content_slides[0]

    # Slayt görevleri: görseller açıkken, sprint toplam SP'sinin %MIN_PCT'inden
    # FAZLA efor alan parent'lar ayrı slayt (rozetli); kalanlar 'Diğer Çalışmalar'.
    # Aksi halde eski davranış (grup/slayt).
    total_sp = sum(s.sp for s in summaries) or 1.0
    min_pct = _min_pct()
    if _visuals_enabled():
        indiv = [s for s in summaries if 100.0 * s.sp / total_sp > min_pct]
        rest = [s for s in summaries if 100.0 * s.sp / total_sp <= min_pct]
        tasks = [("top", [s]) for s in indiv]
        per_dig = _digest_per_slide()
        for i in range(0, len(rest), per_dig):
            tasks.append(("digest", rest[i : i + per_dig]))
        if not tasks:
            tasks = [("top", [])]
    else:
        per = _groups_per_slide()
        tasks = [("top", summaries[i : i + per]) for i in range(0, len(summaries), per)] or [
            ("top", [])
        ]

    # fazladan kalan şablon içerik slaytlarını sldIdLst'ten çıkar
    if len(tasks) < len(content_slides):
        for i in content_idx[len(tasks):]:
            sldIdLst.remove(el_of[i])
        content_slides = content_slides[: len(tasks)]

    # gereğinden az slayt varsa klonla
    used = list(content_slides)
    while len(used) < len(tasks):
        used.append(_clone_slide(prs, src))

    start_no = 1
    for (kind, chunk), slide in zip(tasks, used):
        if kind == "digest":
            _set_slide_title(slide, "Diğer Çalışmalar")
            _fill_digest_slide(slide, chunk)
        else:
            _fill_content_slide(slide, chunk, start_no)
            start_no += len(chunk)

    _set_team_image(prs, team_image_path())
    _reorder_slides(prs)
    _renumber_pages(prs)

    prs.save(out_path)
    return out_path


# ------------------------------------------------------------------ entry point


def _fmt_date(iso: str) -> str:
    if not iso:
        return ""
    try:
        return datetime.fromisoformat(iso.replace("Z", "+00:00")).strftime("%d.%m.%Y")
    except Exception:
        return iso[:10]


def _pretty_sprint(name: str) -> str:
    """'2026_13_Sudo' -> '13. Sprint' (FLO deck konvansiyonu); eşleşmezse aynen döner."""
    m = re.match(r"^\d{4}[_-](\d{1,2})[_-]", name or "")
    if m:
        return f"{int(m.group(1))}. Sprint"
    return name or ""


def _sprint_meta(client, iteration_path: str, team: str = "") -> dict:
    """list_iterations içinden bu iteration'ın adı + tarih aralığı (tarihler attributes altında)."""
    try:
        iters = client.list_iterations(team=team)
    except Exception:
        iters = []
    for it in iters:
        if it.get("path") == iteration_path or it.get("name") == iteration_path:
            attrs = it.get("attributes", {}) or {}
            return {
                "name": it.get("name", ""),
                "startDate": attrs.get("startDate", ""),
                "finishDate": attrs.get("finishDate", ""),
            }
    return {"name": iteration_path.split("\\")[-1], "startDate": "", "finishDate": ""}


# ------------------------------------------------------------------ Azure Analytics (burndown/velocity)

_WI_TYPE_FILTER = "(WorkItemType eq 'Task' or WorkItemType eq 'Bug')"
# Azure "Show Resolved work items as Completed" davranışı: Resolved da 'tamamlanmış' sayılır
_DONE_CATEGORIES = ("Completed", "Resolved")


def _last_completed_sprints(client, team: str, n: int) -> list[dict]:
    """Takımın son n TAMAMLANMIŞ sprinti (kronolojik). Her biri {name,startDate,finishDate}."""
    try:
        iters = client.list_iterations(team=team)
    except Exception:
        return []
    past = [
        {
            "name": it.get("name", ""),
            "startDate": it.get("attributes", {}).get("startDate", ""),
            "finishDate": it.get("attributes", {}).get("finishDate", ""),
        }
        for it in iters
        if it.get("attributes", {}).get("timeFrame") == "past"
        and it.get("attributes", {}).get("startDate")
    ]
    past.sort(key=lambda s: s["startDate"])
    return past[-n:]


def burndown_data(client, iteration_name: str, start_iso: str, finish_iso: str) -> dict | None:
    """Azure Analytics'ten günlük burndown: {labels, remaining, total, ideal}.
    remaining = tamamlanmamış SP; total = toplam kapsam; ideal = doğrusal düşüş."""
    if not (iteration_name and start_iso and finish_iso):
        return None
    s, f = start_iso[:10], finish_iso[:10]
    nm = iteration_name.replace("'", "''")
    q = (
        "WorkItemSnapshot?$apply=filter("
        f"Iteration/IterationName eq '{nm}' and DateValue ge {s}Z and DateValue le {f}Z "
        f"and {_WI_TYPE_FILTER})"
        "/groupby((DateValue,StateCategory),aggregate(Custom_StoryPoints with sum as SP))"
        "&$orderby=DateValue"
    )
    rows = client.analytics_query(q)
    if not rows:
        return None
    from collections import defaultdict

    byday: dict[str, dict] = defaultdict(lambda: {"rem": 0.0, "tot": 0.0})
    for r in rows:
        d = (r.get("DateValue") or "")[:10]
        sp = r.get("SP") or 0.0
        byday[d]["tot"] += sp
        if r.get("StateCategory") not in _DONE_CATEGORIES:  # Completed + Resolved düşülür
            byday[d]["rem"] += sp
    days = sorted(byday)
    if not days:
        return None
    remaining = [round(byday[d]["rem"], 1) for d in days]
    total = [round(byday[d]["tot"], 1) for d in days]
    labels = [f"{d[8:10]}.{d[5:7]}" for d in days]  # dd.mm
    n = len(days)
    start_rem = remaining[0]
    ideal = [round(start_rem * (1 - i / (n - 1)), 1) if n > 1 else 0.0 for i in range(n)]
    return {"labels": labels, "remaining": remaining, "total": total, "ideal": ideal}


def velocity_data(client, team: str, n: int = 6) -> dict | None:
    """Son n tamamlanmış sprint için {sprints, planned, total, done} (Custom_StoryPoints).
    planned = başlangıç snapshot; total = bitiş snapshot; done = şu an Done (StateCategory Completed)."""
    sprints = _last_completed_sprints(client, team, n)
    if not sprints:
        return None
    names, planned, total, done = [], [], [], []
    for sp in sprints:
        nm = sp["name"]
        nm_q = nm.replace("'", "''")
        s, f = sp["startDate"][:10], sp["finishDate"][:10]
        # başlangıç + bitiş snapshot (Planlanan / Toplam)
        snap_q = (
            "WorkItemSnapshot?$apply=filter("
            f"Iteration/IterationName eq '{nm_q}' and (DateValue eq {s}Z or DateValue eq {f}Z) "
            f"and {_WI_TYPE_FILTER})"
            "/groupby((DateValue),aggregate(Custom_StoryPoints with sum as SP))"
        )
        try:
            snap = client.analytics_query(snap_q)
        except Exception:
            snap = []
        p = next((r.get("SP") or 0.0 for r in snap if (r.get("DateValue") or "")[:10] == s), 0.0)
        t = next((r.get("SP") or 0.0 for r in snap if (r.get("DateValue") or "")[:10] == f), 0.0)
        # şu an Done (WorkItems, current state)
        done_q = (
            "WorkItems?$apply=filter("
            f"Iteration/IterationName eq '{nm_q}' "
            "and (StateCategory eq 'Completed' or StateCategory eq 'Resolved') "
            f"and {_WI_TYPE_FILTER})"
            "/aggregate(Custom_StoryPoints with sum as SP)"
        )
        try:
            dq = client.analytics_query(done_q)
            d = (dq[0].get("SP") or 0.0) if dq else 0.0
        except Exception:
            d = 0.0
        names.append(nm)
        planned.append(round(p, 1))
        total.append(round(t, 1))
        done.append(round(d, 1))
    return {"sprints": names, "planned": planned, "total": total, "done": done}


def generate_sprint_report(
    team: str,
    iteration_path: str,
    client=None,
    db=None,
    vector_store=None,
    report_id: str = "",
) -> dict:
    """Sprint için .pptx üret. {file_path, group_count, item_count, ai} döner."""
    if _env("CREW_SPRINT_REPORT", "0") in ("0", "false", "False", ""):
        raise RuntimeError("CREW_SPRINT_REPORT kapalı")

    tmpl = template_path()
    if not os.path.exists(tmpl):
        raise FileNotFoundError(
            f"Sunum şablonu bulunamadı: {tmpl} (CREW_SPRINT_REPORT_TEMPLATE ile ayarla)"
        )

    if client is None:
        from agile_sdlc_crew.tools.azure_devops_base import AzureDevOpsClient

        client = AzureDevOpsClient()
    if db is None:
        from agile_sdlc_crew import db as db  # noqa: PLW0127
    if vector_store is None:
        try:
            from agile_sdlc_crew.tools.vector_store import VectorStore

            vector_store = VectorStore()
        except Exception as e:
            log.warning(f"  VectorStore alınamadı, /repo-decisions atlanıyor: {e}")
            vector_store = None

    items = collect_sprint_items(client, iteration_path)
    groups = group_by_parent(items)

    repo_decisions = {}
    if vector_store is not None:
        try:
            repo_decisions = vector_store.repo_decision_records()
        except Exception as e:
            log.warning(f"  repo-decisions okunamadı: {e}")

    summaries: list[GroupSummary] = []
    for g in groups:
        contexts = {
            it.id: _stored_context_for_item(it, db, repo_decisions, client) for it in g.items
        }
        summaries.append(summarize_group(g, contexts))

    # sprint toplam SP'si içindeki oranlar
    total_sp = sum(s.sp for s in summaries)
    if total_sp > 0:
        for s in summaries:
            s.sp_pct = round(100.0 * s.sp / total_sp)

    meta = _sprint_meta(client, iteration_path, team)
    sprint_name = meta.get("name", iteration_path)
    date_range = ""
    if meta.get("startDate") or meta.get("finishDate"):
        date_range = f"{_fmt_date(meta.get('startDate',''))} – {_fmt_date(meta.get('finishDate',''))}"
    title_line = _env("CREW_SPRINT_REPORT_TITLE") or team or "Sprint Raporu"

    # Azure Analytics: burndown (mevcut sprint) + velocity (son N sprint)
    burndown = velocity = None
    if _azure_charts_enabled():
        try:
            burndown = burndown_data(client, sprint_name, meta.get("startDate", ""), meta.get("finishDate", ""))
        except Exception as e:
            log.warning(f"  Burndown verisi alınamadı: {e}")
        try:
            velocity = velocity_data(client, team, _velocity_sprints())
        except Exception as e:
            log.warning(f"  Velocity verisi alınamadı: {e}")

    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", f"{team}_{sprint_name}").strip("_")
    fname = f"sprint_report_{safe}_{report_id or 'out'}.pptx"
    out_path = os.path.join(output_dir(), fname)

    build_pptx(
        title_line, _pretty_sprint(sprint_name), date_range, summaries, tmpl, out_path,
        burndown=burndown, velocity=velocity,
    )

    return {
        "file_path": out_path,
        "file_name": fname,
        "group_count": len(groups),
        "item_count": len(items),
        "ai": _ai_enabled(),
        "sprint_name": sprint_name,
    }
